from pptx import Presentation
from bs4 import BeautifulSoup
import os
import olefile


# 폰트 추출 함수 (.pptx, .hml 모두 지원)
def extract_fonts(filepath):
    fonts = set()

    # PPTX 파일 처리
    if filepath.endswith(".pptx"):
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                fonts.add(run.font.name)

    # HWP(HML 변환된 파일) 처리
    elif filepath.endswith(".hml"):
        with open(filepath, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "xml")
            for font_tag in soup.find_all("FONTFACE"):
                font_name = font_tag.get("FACE")
                if font_name:
                    fonts.add(font_name)

    return fonts

# 누락된 폰트 찾기
def find_missing_fonts(fonts):
    installed_fonts = set(os.listdir("C:/Windows/Fonts"))  # 윈도우 폰트 폴더 기준
    missing = [f for f in fonts if f not in ''.join(installed_fonts)]
    return missing


import olefile

def extract_fonts_from_hwp(filepath):
    fonts = set()

    if filepath.endswith('.hwp'):
        ole = olefile.OleFileIO(filepath)

        # 한글 파일 내에서 'PrvText' 스트림은 문서 텍스트가 들어있는 부분
        if ole.exists('PrvText'):
            stream = ole.openstream('PrvText')
            content = stream.read().decode('utf-16', errors='ignore')
            
            # 한글 폰트 이름은 일반적으로 텍스트에 포함되므로, 몇 가지 폰트 패턴을 탐지
            font_candidates = ["바탕", "굴림", "돋움", "맑은 고딕", "HY", "한컴", "함초롬", "Arial", "Times", "Gothic", "Sans"]
            for name in font_candidates:
                if name in content:
                    fonts.add(name)

        ole.close()

    return fonts
